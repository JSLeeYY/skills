from pathlib import Path
import shutil

from pptx import Presentation
from pptx.oxml.ns import qn


SOURCE_PPT = Path(
    r"D:\DevelopmentLocation\agent skill\skills\word_project\2026\2026-04-03\汇报材料\监造信息化系统甲方汇报-一期建设方案-V4.pptx"
)
TARGET_PPT = SOURCE_PPT.with_name(
    SOURCE_PPT.name.replace("V4", "V5")
)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide = slide_id_list[index]
    rel_id = slide.get(qn("r:id"))
    prs.part.drop_rel(rel_id)
    del slide_id_list[index]


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def replace_first_prefix(slide, prefix: str, new_text: str) -> bool:
    for shape in text_shapes(slide):
        text = shape.text.strip()
        if text.startswith(prefix):
            shape.text = new_text
            return True
    return False


def remove_prefix_shapes(slide, prefixes) -> None:
    prefixes = tuple(prefixes)
    for shape in list(text_shapes(slide)):
        text = shape.text.strip()
        if text.startswith(prefixes):
            remove_shape(shape)


def replace_left_bottom_text(slide, new_text: str) -> bool:
    for shape in text_shapes(slide):
        text = shape.text.strip()
        if (
            text.startswith("•")
            and shape.left < 2_000_000
            and shape.top > 3_500_000
        ):
            shape.text = new_text
            return True
    return False


def strip_lines_with_prefix(shape, prefixes) -> None:
    prefixes = tuple(prefixes)
    lines = [line.strip() for line in shape.text.splitlines()]
    kept = [line for line in lines if line and not line.startswith(prefixes)]
    if kept:
        shape.text = "\n".join(kept)


def sanitize_cover(slide) -> None:
    replace_first_prefix(slide, "一期建设方案 V4", "一期建设方案 V5")
    replace_first_prefix(slide, "监造信息化系统甲方汇报", "监造信息化系统甲方汇报")
    replace_first_prefix(
        slide,
        "口径基准：",
        "聚焦一期建设范围、核心流程和模块关系。",
    )
    replace_first_prefix(slide, "本版重构重点", "汇报范围")
    replace_first_prefix(
        slide,
        "• 不沿用 V3",
        "• 业务全景与模块关系\n• 主链流程与支撑模块\n• 一期建设边界",
    )
    replace_first_prefix(slide, "汇报对象：", "版本：甲方留存版")


def sanitize_overview(slide) -> None:
    replace_first_prefix(
        slide,
        "按 5 个阶段理解整个系统：",
        "按 5 个阶段展示系统范围，一期聚焦主链流程和关键支撑能力。",
    )
    remove_prefix_shapes(slide, ["本页核心："])
    replace_first_prefix(
        slide,
        "模块数量口径：",
        "模块口径：主链 15 个，一期支撑 10 个，二期增强 2 个。",
    )
    for shape in text_shapes(slide):
        if "价值：" in shape.text:
            strip_lines_with_prefix(shape, ["价值："])


def sanitize_relationship(slide) -> None:
    replace_first_prefix(
        slide,
        "主链按业务先后关系展开，问题处理和支撑模块以挂载方式接入，不再强行画成单线。",
        "主链按业务先后展开，支撑模块按挂载方式接入。",
    )
    remove_prefix_shapes(slide, ["本页核心："])
    replace_first_prefix(
        slide,
        "支撑模块挂载逻辑：",
        "支撑模块：统计分析、风险预警、可视化、工日、二维码、签名、OCR、视频、付款、培训。",
    )


def sanitize_core_module(slide) -> None:
    replace_first_prefix(
        slide,
        "一页讲清一个核心模块，流程图、规则、按钮、角色和上下游一次展开。",
        "核心模块流程与职责。",
    )
    remove_prefix_shapes(slide, ["本页核心："])
    replace_first_prefix(slide, "关键业务规则 / 字段 / 限制", "模块说明")
    replace_left_bottom_text(
        slide,
        "说明：本页保留流程、功能按钮、参与角色和上下游关系，详细规则按实施阶段需求确认。",
    )


def sanitize_support_module(slide) -> None:
    replace_first_prefix(
        slide,
        "支撑模块不替代主链，但必须挂到关键节点上真正发生作用。",
        "支撑模块挂载到主链关键节点。",
    )
    remove_prefix_shapes(slide, ["本页核心："])
    replace_first_prefix(slide, "落地要点", "模块说明")
    replace_left_bottom_text(
        slide,
        "说明：本模块按一期范围挂载到对应业务节点，详细配置在实施阶段确认。",
    )


def build_client_v5() -> Path:
    shutil.copy2(SOURCE_PPT, TARGET_PPT)
    prs = Presentation(TARGET_PPT)

    while len(prs.slides) > 28:
        delete_slide(prs, 28)

    sanitize_cover(prs.slides[0])
    sanitize_overview(prs.slides[1])
    sanitize_relationship(prs.slides[2])

    for slide_index in range(3, 18):
        sanitize_core_module(prs.slides[slide_index])

    for slide_index in range(18, 28):
        sanitize_support_module(prs.slides[slide_index])

    prs.save(TARGET_PPT)
    return TARGET_PPT


if __name__ == "__main__":
    output_path = build_client_v5()
    print(output_path)
