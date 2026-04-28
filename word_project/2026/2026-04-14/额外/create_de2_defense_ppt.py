from __future__ import annotations

import re
import shutil
from dataclasses import dataclass
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent

COLOR_BLUE = "327DF9"
COLOR_GREEN = "7CC636"
COLOR_ORANGE = "FFA436"
COLOR_CYAN = "54C3F1"
COLOR_PURPLE = "A867F7"
COLOR_DARK = "373C53"
COLOR_TEXT = "2F3A4F"
COLOR_MUTED = "667085"
COLOR_BG = "F6F9FC"
COLOR_WHITE = "FFFFFF"


@dataclass
class ProjectRecord:
    project_output: float
    name: str
    modules: str
    pages: int


@dataclass
class ApplicantData:
    name: str
    target_level: str
    projects: list[ProjectRecord]
    quality: str
    feedback: str
    certification: str
    capability: str
    process: str
    total_output: float
    total_pages: int
    mentee_names: list[str]


@dataclass
class StandardData:
    project_output: str
    showcase: str
    issue_quality: str
    knowledge_feedback: str
    mentoring: str
    sharing: str
    certification: str
    domain_config: str
    logic_service: str
    problem_solving: str
    data_development: str
    process: str
    config_standard: str


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def lighten(hex_value: str, ratio: float) -> str:
    base = [int(hex_value[i : i + 2], 16) for i in range(0, 6, 2)]
    lifted = [int(channel + (255 - channel) * ratio) for channel in base]
    return "".join(f"{value:02X}" for value in lifted)


def clean_text(text: str | None) -> str:
    if not text:
        return ""
    normalized = str(text).replace("\r", "\n")
    normalized = re.sub(r"[ \t]+", " ", normalized)
    normalized = re.sub(r"\n{2,}", "\n", normalized)
    return normalized.strip()


def split_lines(text: str | None) -> list[str]:
    return [line.strip(" ;；") for line in clean_text(text).split("\n") if line.strip(" ;；")]


def short_text(text: str, limit: int) -> str:
    text = clean_text(text)
    return text if len(text) <= limit else f"{text[: limit - 1]}…"


def find_single_file(suffix: str, *tokens: str) -> Path:
    candidates = []
    for path in ROOT.iterdir():
        if not path.is_file() or path.suffix.lower() != suffix.lower():
            continue
        if all(token in path.name for token in tokens):
            candidates.append(path)
    if len(candidates) != 1:
        raise FileNotFoundError(f"Unable to locate a unique file for {suffix} with tokens {tokens}")
    return candidates[0]


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide = slide_id_list[index]
    rel_id = slide.get(qn("r:id"))
    prs.part.drop_rel(rel_id)
    del slide_id_list[index]


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        delete_slide(prs, 0)


def get_placeholder(slide, idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    raise KeyError(f"Placeholder idx={idx} not found")


def apply_run_style(run, *, size: float, color: str = COLOR_TEXT, bold: bool = False, font_name: str = "Microsoft YaHei") -> None:
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.name = font_name
    font.color.rgb = rgb(color)


def write_text_frame(
    text_frame,
    paragraphs: list[dict],
    *,
    margin: float = 0.08,
    wrap: bool = True,
    vertical: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = wrap
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = vertical
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)

    for index, spec in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = spec.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(spec.get("space_after", 0))
        paragraph.space_before = Pt(spec.get("space_before", 0))
        paragraph.level = spec.get("level", 0)
        run = paragraph.add_run()
        run.text = spec["text"]
        apply_run_style(
            run,
            size=spec.get("size", 12),
            color=spec.get("color", COLOR_TEXT),
            bold=spec.get("bold", False),
            font_name=spec.get("font_name", "Microsoft YaHei"),
        )


def set_placeholder_text(shape, text: str, *, size: float, color: str = COLOR_TEXT, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    write_text_frame(
        shape.text_frame,
        [{"text": text, "size": size, "color": color, "bold": bold, "align": align}],
        margin=0.0,
    )


def add_shape(slide, shape_type, left: float, top: float, width: float, height: float, *, fill: str, line: str, radius: bool = False):
    actual_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(actual_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.2)
    return shape


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[dict],
    *,
    margin: float = 0.1,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    write_text_frame(textbox.text_frame, paragraphs, margin=margin)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    title: str,
    body_lines: list[str],
    accent: str,
    title_size: float = 15,
    body_size: float = 10.5,
    footer: str | None = None,
    fill: str | None = None,
) -> None:
    fill_color = fill or lighten(accent, 0.90)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.color.rgb = rgb(accent)
    shape.line.width = Pt(1.2)

    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.16),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = rgb(accent)
    strip.line.color.rgb = rgb(accent)

    paragraphs = [{"text": title, "size": title_size, "bold": True, "color": accent, "space_after": 3}]
    paragraphs.extend({"text": f"• {line}", "size": body_size, "color": COLOR_TEXT, "space_after": 2} for line in body_lines if line)
    if footer:
        paragraphs.append({"text": footer, "size": body_size, "color": COLOR_MUTED, "space_before": 4})

    add_textbox(slide, left + 0.12, top + 0.18, width - 0.24, height - 0.26, paragraphs)


def add_metric_card(slide, left: float, top: float, width: float, height: float, value: str, label: str, accent: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(lighten(accent, 0.88))
    shape.line.color.rgb = rgb(accent)
    shape.line.width = Pt(1.2)
    paragraphs = [
        {"text": value, "size": 24, "bold": True, "color": accent, "align": PP_ALIGN.CENTER, "space_after": 2},
        {"text": label, "size": 11.5, "color": COLOR_DARK, "align": PP_ALIGN.CENTER},
    ]
    write_text_frame(shape.text_frame, paragraphs, margin=0.08, vertical=MSO_VERTICAL_ANCHOR.MIDDLE)


def style_table_cell(cell, text: str, *, fill: str, color: str, size: float, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)
    cell.margin_left = Pt(8)
    cell.margin_right = Pt(8)
    cell.margin_top = Pt(6)
    cell.margin_bottom = Pt(6)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    write_text_frame(
        cell.text_frame,
        [{"text": text, "size": size, "color": color, "bold": bold, "align": align}],
        margin=0.0,
    )


def style_content_title(slide, title: str) -> None:
    title_shape = get_placeholder(slide, 0)
    set_placeholder_text(title_shape, title, size=24, color=COLOR_DARK, bold=True)


def parse_applicant_data(path: Path) -> ApplicantData:
    wb = load_workbook(path, data_only=True)
    ws = wb.active

    applicant_name = clean_text(ws.cell(1, 3).value)
    target_level = clean_text(ws.cell(1, 7).value)

    quality = clean_text(ws.cell(3, 5).value)
    feedback = clean_text(ws.cell(3, 6).value)
    certification = clean_text(ws.cell(3, 7).value)
    capability = clean_text(ws.cell(3, 8).value)
    process = clean_text(ws.cell(3, 9).value)

    projects: list[ProjectRecord] = []
    for row in range(3, ws.max_row + 1):
        project_name = clean_text(ws.cell(row, 2).value)
        pages_text = clean_text(ws.cell(row, 4).value)
        if not project_name:
            continue
        project_output_text = clean_text(ws.cell(row, 1).value)
        pages = int(float(pages_text)) if pages_text else 0
        project_output = float(project_output_text) if project_output_text else 0.0
        projects.append(
            ProjectRecord(
                project_output=project_output,
                name=project_name,
                modules=clean_text(ws.cell(row, 3).value),
                pages=pages,
            )
        )

    total_output = round(sum(project.project_output for project in projects), 1)
    total_pages = sum(project.pages for project in projects)

    mentee_names: list[str] = []
    feedback_lines = split_lines(feedback)
    if feedback_lines:
        match = re.search(r"辅导(.+?)掌握新技能", feedback_lines[0])
        if match:
            names = re.split(r"[，,、/ ]+", match.group(1))
            mentee_names = [name for name in names if name]

    return ApplicantData(
        name=applicant_name,
        target_level=target_level,
        projects=projects,
        quality=quality,
        feedback=feedback,
        certification=certification,
        capability=capability,
        process=process,
        total_output=total_output,
        total_pages=total_pages,
        mentee_names=mentee_names,
    )


def parse_standard_data(path: Path) -> StandardData:
    wb = load_workbook(path, data_only=True)
    ws = wb[wb.sheetnames[0]]

    last_dim_1 = ""
    last_dim_2 = ""
    de2_lookup: dict[tuple[str, str], str] = {}

    for row in ws.iter_rows(values_only=True):
        dim_1 = clean_text(row[0]) or last_dim_1
        dim_2 = clean_text(row[1]) or last_dim_2
        dim_3 = clean_text(row[2])
        de2_value = clean_text(row[5]) if len(row) > 5 else ""
        if clean_text(row[0]):
            last_dim_1 = dim_1
        if clean_text(row[1]):
            last_dim_2 = dim_2
        if dim_2 or dim_3:
            de2_lookup[(dim_2, dim_3)] = de2_value

    return StandardData(
        project_output=de2_lookup[("项目产出", "初级项目")],
        showcase=de2_lookup[("项目质量", "模块showcase")],
        issue_quality=de2_lookup[("项目质量", "问题单")],
        knowledge_feedback=de2_lookup[("专业回馈", "知识经验")],
        mentoring=de2_lookup[("专业回馈", "技能传递")],
        sharing=de2_lookup[("专业回馈", "项目分享")],
        certification=de2_lookup[("专业认证", "证书要求")],
        domain_config=de2_lookup[("配置能力", "三域配置能力")],
        logic_service=de2_lookup[("配置能力", "逻辑与服务")],
        problem_solving=de2_lookup[("配置能力", "问题定位能力")],
        data_development=de2_lookup[("配置能力", "数据开发能力")],
        process=de2_lookup[("流程规范", "过程管理")],
        config_standard=de2_lookup[("流程规范", "配置规范")],
    )


def build_contents_slide(slide) -> None:
    titles = [
        "个人概览与申报结论",
        "项目产出与交付质量",
        "专业回馈与能力对标",
        "总结与后续计划",
    ]
    accents = [COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE, COLOR_CYAN]
    top_positions = [1.92, 3.00, 4.08, 5.16]

    note_box = slide.shapes.add_textbox(Inches(0.95), Inches(0.8), Inches(4.0), Inches(0.5))
    write_text_frame(
        note_box.text_frame,
        [{"text": "DE2 申报答辩提纲", "size": 24, "bold": True, "color": COLOR_DARK}],
        margin=0.0,
    )

    for index, (title, accent, top) in enumerate(zip(titles, accents, top_positions), start=1):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(5.85),
            Inches(top),
            Inches(6.2),
            Inches(0.62),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(lighten(accent, 0.94))
        card.line.color.rgb = rgb(accent)
        card.line.width = Pt(1.2)

        number = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(5.45),
            Inches(top - 0.06),
            Inches(0.62),
            Inches(0.62),
        )
        number.fill.solid()
        number.fill.fore_color.rgb = rgb(accent)
        number.line.color.rgb = rgb(accent)
        write_text_frame(
            number.text_frame,
            [{"text": f"{index:02d}", "size": 17, "bold": True, "color": COLOR_WHITE, "align": PP_ALIGN.CENTER}],
            margin=0.0,
            vertical=MSO_VERTICAL_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            6.25,
            top + 0.08,
            5.35,
            0.32,
            [{"text": title, "size": 18, "bold": True, "color": COLOR_DARK}],
            margin=0.0,
        )


def build_section_slide(slide, title: str, subtitle: str) -> None:
    main_shape = get_placeholder(slide, 12)
    sub_shape = get_placeholder(slide, 14)
    set_placeholder_text(main_shape, title, size=28, color=COLOR_DARK, bold=True)
    set_placeholder_text(sub_shape, subtitle, size=14, color=COLOR_MUTED)


def build_cover(slide, applicant: ApplicantData) -> None:
    title_shape = get_placeholder(slide, 12)
    sub_shape = get_placeholder(slide, 14)
    footer_shape = get_placeholder(slide, 13)

    set_placeholder_text(title_shape, "配置开发专业任职资格答辩", size=28, color=COLOR_DARK, bold=True)
    set_placeholder_text(sub_shape, f"{applicant.target_level} 申报答辩", size=18, color=COLOR_BLUE, bold=True)
    set_placeholder_text(footer_shape, f"汇报人：{applicant.name}    2026.04", size=11.5, color=COLOR_MUTED)


def build_overview_slide(slide, applicant: ApplicantData) -> None:
    style_content_title(slide, "个人概览：以独立交付能力支撑本次 DE2 申报")

    metric_top = 1.35
    metric_lefts = [0.95, 3.52, 6.09, 8.66]
    metric_width = 2.32
    metric_height = 1.05
    metric_values = [
        (f"{len(applicant.projects)}个", "参与项目", COLOR_BLUE),
        (f"{applicant.total_output:.1f}", "项目产出", COLOR_GREEN),
        (f"{applicant.total_pages}页", "输出页面", COLOR_ORANGE),
        (f"{max(len(applicant.mentee_names), 4)}位", "辅导同事", COLOR_CYAN),
    ]
    for left, (value, label, accent) in zip(metric_lefts, metric_values):
        add_metric_card(slide, left, metric_top, metric_width, metric_height, value, label, accent)

    add_card(
        slide,
        0.95,
        2.62,
        5.4,
        3.25,
        title="个人定位",
        body_lines=[
            f"当前以 {applicant.target_level} 为申报目标，已具备稳定的独立配置开发能力。",
            "累计参与招投标、预算、制造、大屏、环保、安全、采购合同等多类业务场景。",
            "能够在多项目并行环境下，独立完成模块配置、联调、问题定位与交付支撑。",
        ],
        accent=COLOR_BLUE,
        body_size=11.5,
    )

    add_card(
        slide,
        6.55,
        2.62,
        5.5,
        3.25,
        title="申报结论",
        body_lines=[
            "项目产出、质量表现、配置能力与流程协同，已经形成较完整的 DE2 申报基础。",
            "当前已通过 SDE 认证，并围绕更高标准继续强化分享沉淀与复用输出。",
            "本次答辩的核心目标，是证明我已具备 DE2 所要求的独立交付与稳定质量能力。",
        ],
        accent=COLOR_GREEN,
        body_size=11.5,
    )


def build_standard_slide(slide, applicant: ApplicantData, standards: StandardData) -> None:
    style_content_title(slide, "DE2 标准对标：围绕关键维度说明本次申报基础")

    rows = [
        ("项目产出", "不少于 6 个项目产出", f"按申报表口径累计 {applicant.total_output:.1f} 个项目产出，覆盖 {len(applicant.projects)} 个项目"),
        ("项目质量", "独立 showcase 通过率 90%，问题单占比小于 3", "所有项目均个人独立 showcase，通过率 90%；近 1 个月无问题单打回"),
        ("专业回馈", "总结/改进建议、辅导与项目分享持续输出", "已辅导 4 位同事、完成内部分享与项目模块分享，持续做设计/接口总结"),
        ("专业认证", short_text(standards.certification, 20), f"当前已获得 {applicant.certification}，以此作为 {applicant.target_level} 申报基础"),
        ("配置能力", "熟练掌握平台功能并能独立完成配置开发任务", "熟练使用逻辑控制、表单列表、业务流、交换机，并可独立完成连接器/接口配置"),
        ("流程规范", short_text(standards.process, 26), clean_text(applicant.process)),
    ]

    table_shape = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.95), Inches(1.45), Inches(11.15), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(1.45)
    table.columns[1].width = Inches(3.75)
    table.columns[2].width = Inches(5.95)

    headers = ["维度", "DE2 标准要点", "我的支撑说明"]
    for column, header in enumerate(headers):
        style_table_cell(table.cell(0, column), header, fill=COLOR_DARK, color=COLOR_WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)

    row_fills = [lighten(COLOR_BLUE, 0.93), COLOR_WHITE]
    for row_index, (dimension, standard_text, evidence_text) in enumerate(rows, start=1):
        fill_color = row_fills[row_index % 2]
        style_table_cell(table.cell(row_index, 0), dimension, fill=lighten(COLOR_BLUE, 0.84), color=COLOR_DARK, size=11.5, bold=True, align=PP_ALIGN.CENTER)
        style_table_cell(table.cell(row_index, 1), standard_text, fill=fill_color, color=COLOR_TEXT, size=10.2)
        style_table_cell(table.cell(row_index, 2), evidence_text, fill=fill_color, color=COLOR_TEXT, size=10.2)

    add_card(
        slide,
        0.95,
        6.02,
        11.15,
        0.82,
        title="阶段判断",
        body_lines=[
            "项目产出、质量、配置能力与流程协同已具备明确支撑。",
            "认证与专业回馈在本次答辩中重点展示基础与后续强化方向。",
        ],
        accent=COLOR_CYAN,
        title_size=13,
        body_size=10.5,
        fill=lighten(COLOR_CYAN, 0.92),
    )


def build_project_output_slide(slide, applicant: ApplicantData) -> None:
    style_content_title(slide, "项目产出：覆盖 9 个项目，形成 10.9 个项目产出与 367 页交付")

    projects = sorted(applicant.projects, key=lambda item: item.pages, reverse=True)

    table_shape = slide.shapes.add_table(len(projects) + 1, 3, Inches(0.95), Inches(1.45), Inches(6.2), Inches(4.95))
    table = table_shape.table
    table.columns[0].width = Inches(4.55)
    table.columns[1].width = Inches(0.8)
    table.columns[2].width = Inches(0.85)

    for column, header in enumerate(["项目名称", "产出", "页数"]):
        style_table_cell(table.cell(0, column), header, fill=COLOR_DARK, color=COLOR_WHITE, size=11.5, bold=True, align=PP_ALIGN.CENTER)

    for row_index, project in enumerate(projects, start=1):
        fill_color = COLOR_WHITE if row_index % 2 else lighten(COLOR_BLUE, 0.95)
        style_table_cell(table.cell(row_index, 0), short_text(project.name, 28), fill=fill_color, color=COLOR_TEXT, size=9.6)
        style_table_cell(table.cell(row_index, 1), f"{project.project_output:g}", fill=fill_color, color=COLOR_TEXT, size=10.2, align=PP_ALIGN.CENTER)
        style_table_cell(table.cell(row_index, 2), str(project.pages), fill=fill_color, color=COLOR_TEXT, size=10.2, align=PP_ALIGN.CENTER)

    add_card(
        slide,
        7.35,
        1.45,
        4.75,
        0.78,
        title="页面产出 Top5",
        body_lines=[],
        accent=COLOR_GREEN,
        title_size=13.5,
        body_size=10.5,
        fill=lighten(COLOR_GREEN, 0.93),
    )

    chart_items = projects[:5]
    max_pages = max(project.pages for project in chart_items)
    bar_top = 2.28
    for index, project in enumerate(chart_items):
        y = bar_top + index * 0.72
        accent = [COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE, COLOR_CYAN, COLOR_PURPLE][index]
        add_textbox(
            slide,
            7.35,
            y - 0.02,
            2.1,
            0.3,
            [{"text": short_text(project.name, 12), "size": 10.2, "color": COLOR_TEXT}],
            margin=0.0,
        )
        bar_width = 2.15 * project.pages / max_pages + 0.65
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            9.5,
            y,
            bar_width,
            0.23,
            fill=lighten(accent, 0.18),
            line=accent,
            radius=True,
        )
        add_textbox(
            slide,
            11.75,
            y - 0.06,
            0.35,
            0.3,
            [{"text": str(project.pages), "size": 10.5, "color": accent, "bold": True, "align": PP_ALIGN.RIGHT}],
            margin=0.0,
        )

    add_card(
        slide,
        7.35,
        5.2,
        4.75,
        1.18,
        title="业务覆盖",
        body_lines=[
            "覆盖招投标、预算管理、制造平台、大屏驾驶仓、环保分析、安全运营、采购合同等场景。",
            "说明我的交付能力具备较强的迁移和业务适配能力，而非局限于单一模块。",
        ],
        accent=COLOR_ORANGE,
        title_size=13.5,
        body_size=10.1,
        fill=lighten(COLOR_ORANGE, 0.93),
    )


def build_project_quality_slide(slide, applicant: ApplicantData) -> None:
    style_content_title(slide, "代表项目与交付质量：以关键项目说明独立交付与质量表现")

    representative_projects = []
    for project_name in [
        "海外集团电子招投标平台二期项目",
        "长智院宜宾市合同管理业务系统项目/公司级数字化制造平台一期项目",
        "南京炼油厂",
    ]:
        for project in applicant.projects:
            if project.name == project_name:
                representative_projects.append(project)
                break

    card_positions = [0.95, 4.03, 7.11]
    accents = [COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE]
    for left, accent, project in zip(card_positions, accents, representative_projects):
        add_card(
            slide,
            left,
            1.45,
            2.88,
            3.45,
            title=short_text(project.name, 14),
            body_lines=[
                f"项目产出：{project.project_output:g}    输出页面：{project.pages} 页",
                short_text(project.modules or "围绕核心业务模块完成配置开发与联调交付", 54),
                "体现了我在复杂业务链路下的模块配置与交付推进能力。",
            ],
            accent=accent,
            title_size=13.5,
            body_size=9.6,
        )

    add_card(
        slide,
        0.95,
        5.1,
        5.55,
        1.62,
        title="交付质量表现",
        body_lines=[
            "所有项目均为个人独立配置功能 showcase，通过率 90%。",
            "近 2 个月内无配置规范类、自检工具检查问题单。",
            "功能及样式问题单与标准单位产出件总数比小于 3，近 1 个月无问题单打回。",
        ],
        accent=COLOR_CYAN,
        title_size=13.5,
        body_size=10.2,
        fill=lighten(COLOR_CYAN, 0.93),
    )

    add_card(
        slide,
        6.65,
        5.1,
        5.45,
        1.62,
        title="独立交付能力",
        body_lines=[
            "熟练承担模块配置、联调、问题定位、换包升级等完整交付动作。",
            "能够独立与 AE、研发、IE 按流程标准协同，确保产出按时按质完成。",
            "多项目并行时仍能保持交付节奏和质量稳定性。",
        ],
        accent=COLOR_GREEN,
        title_size=13.5,
        body_size=10.2,
        fill=lighten(COLOR_GREEN, 0.93),
    )


def build_feedback_slide(slide, applicant: ApplicantData) -> None:
    style_content_title(slide, "专业回馈：通过带教、分享与文档沉淀反哺团队")

    top_banner = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 1.35, 11.15, 0.78, fill=lighten(COLOR_BLUE, 0.90), line=COLOR_BLUE, radius=True)
    write_text_frame(
        top_banner.text_frame,
        [
            {"text": f"已辅导 {max(len(applicant.mentee_names), 4)} 位同事掌握新技能，并在项目中持续输出分享与文档沉淀。", "size": 16, "bold": True, "color": COLOR_DARK, "align": PP_ALIGN.CENTER}
        ],
        margin=0.05,
        vertical=MSO_VERTICAL_ANCHOR.MIDDLE,
    )

    add_card(
        slide,
        0.95,
        2.42,
        3.45,
        3.6,
        title="技能传递",
        body_lines=[
            "辅导林思含、沈雨、颜雨欣、李亚娇掌握新技能，以支撑项目任务完成。",
            "带教重点不是单次答疑，而是帮助同事形成可落地的配置能力。",
        ],
        accent=COLOR_GREEN,
        title_size=14.5,
        body_size=10.6,
    )

    add_card(
        slide,
        4.58,
        2.42,
        3.45,
        3.6,
        title="项目分享",
        body_lines=[
            "完成新员工内部分享，帮助新人更快理解平台与交付方式。",
            "围绕 812 项目中的密级日志板块做过专题分享，沉淀模块经验。",
            "后续将继续增加跨项目分享频次，扩大个人影响范围。",
        ],
        accent=COLOR_ORANGE,
        title_size=14.5,
        body_size=10.6,
    )

    add_card(
        slide,
        8.21,
        2.42,
        3.89,
        3.6,
        title="文档沉淀",
        body_lines=[
            "每个项目都围绕设计文档优化与接口文档做总结。",
            "在项目推进中同步积累问题处理经验和更优配置路径。",
            "后续将把这些经验进一步沉淀为更标准化、可复用的资产。",
        ],
        accent=COLOR_CYAN,
        title_size=14.5,
        body_size=10.6,
    )

    add_textbox(
        slide,
        0.95,
        6.22,
        11.15,
        0.35,
        [{"text": "专业回馈已经形成日常习惯，下一步重点从“项目经验”提升到“组件化、模板化、跨项目复用资产”。", "size": 11, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}],
        margin=0.0,
    )


def build_capability_slide(slide, applicant: ApplicantData, standards: StandardData) -> None:
    style_content_title(slide, "关键能力对标：认证、配置、问题定位与流程协同")

    certification_panel = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 1.35, 11.15, 0.9, fill=lighten(COLOR_DARK, 0.90), line=COLOR_DARK, radius=True)
    write_text_frame(
        certification_panel.text_frame,
        [
            {"text": f"当前专业认证：{applicant.certification}", "size": 16, "bold": True, "color": COLOR_DARK, "align": PP_ALIGN.CENTER, "space_after": 2},
            {"text": f"DE2 标准证书要求：{standards.certification}", "size": 11.5, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER},
        ],
        margin=0.05,
        vertical=MSO_VERTICAL_ANCHOR.MIDDLE,
    )

    add_card(
        slide,
        0.95,
        2.55,
        5.35,
        1.7,
        title="三域配置与逻辑服务",
        body_lines=[
            "熟练沟通 4 层结构与 3 元组框架。",
            "熟练使用逻辑控制、表单列表、业务流、交换机等核心功能。",
            "逻辑控制模块使用更为深入，可独立完成相关配置开发任务。",
        ],
        accent=COLOR_BLUE,
        title_size=14,
        body_size=10.2,
    )
    add_card(
        slide,
        6.75,
        2.55,
        5.35,
        1.7,
        title="问题定位与实施支撑",
        body_lines=[
            "熟练通过日志和后台定位问题，并掌握换包升级等实施流程。",
            "能够在交付现场快速完成一般问题分析与处理闭环。",
        ],
        accent=COLOR_GREEN,
        title_size=14,
        body_size=10.2,
    )
    add_card(
        slide,
        0.95,
        4.5,
        5.35,
        1.7,
        title="数据开发与接口联调",
        body_lines=[
            "可独立完成交换机、连接器、接口等配置开发工作。",
            "对 JSON 等结构化数据处理具备独立实践能力。",
            "能力覆盖了 DE2 对数据开发能力的核心要求方向。",
        ],
        accent=COLOR_ORANGE,
        title_size=14,
        body_size=10.2,
    )
    add_card(
        slide,
        6.75,
        4.5,
        5.35,
        1.7,
        title="流程规范与跨角色协同",
        body_lines=[
            clean_text(applicant.process),
            short_text(standards.config_standard, 32),
        ],
        accent=COLOR_CYAN,
        title_size=14,
        body_size=10.2,
    )

    add_textbox(
        slide,
        0.95,
        6.42,
        11.15,
        0.28,
        [{"text": "能力侧重点已经从单点配置执行，逐步转向可独立交付、可排查问题、可带动协同的 DE2 能力结构。", "size": 11, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}],
        margin=0.0,
    )


def build_summary_slide(slide, applicant: ApplicantData) -> None:
    style_content_title(slide, "总结与后续计划")

    add_card(
        slide,
        0.95,
        1.45,
        5.45,
        4.5,
        title="本次申报总结",
        body_lines=[
            f"以 {applicant.total_output:.1f} 个项目产出、{applicant.total_pages} 页交付，支撑本次 {applicant.target_level} 申报答辩。",
            "已经具备独立完成配置开发任务、控制交付质量、完成跨角色协同的基础能力。",
            "在项目实践中逐步形成带教、分享与文档沉淀习惯，具备持续成长的正向循环。",
        ],
        accent=COLOR_BLUE,
        title_size=16,
        body_size=11.2,
    )

    add_card(
        slide,
        6.65,
        1.45,
        5.45,
        4.5,
        title="后续强化方向",
        body_lines=[
            "继续把项目中的总结沉淀为可复用组件、模板和标准化方法。",
            "增加跨项目模块分享和专题输出，扩大专业回馈的广度与深度。",
            "进一步强化复杂场景下的数据开发、服务编排与问题定位能力。",
            "从独立完成模块交付，逐步向承担更复杂模块目标和影响范围迈进。",
        ],
        accent=COLOR_GREEN,
        title_size=16,
        body_size=11.2,
    )

    summary_box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 6.15, 11.15, 0.55, fill=lighten(COLOR_ORANGE, 0.92), line=COLOR_ORANGE, radius=True)
    write_text_frame(
        summary_box.text_frame,
        [
            {
                "text": f"申请结论：我已具备 {applicant.target_level} 申报所需的核心基础，期待通过本次答辩完成能力认证。",
                "size": 14,
                "bold": True,
                "color": COLOR_DARK,
                "align": PP_ALIGN.CENTER,
            }
        ],
        margin=0.05,
        vertical=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def build_end_slide(slide, applicant: ApplicantData) -> None:
    add_textbox(
        slide,
        3.25,
        2.28,
        6.9,
        0.9,
        [{"text": "谢谢聆听", "size": 30, "bold": True, "color": COLOR_DARK, "align": PP_ALIGN.CENTER}],
        margin=0.0,
    )
    add_textbox(
        slide,
        3.25,
        3.15,
        6.9,
        0.42,
        [{"text": "请各位评审指正", "size": 14, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}],
        margin=0.0,
    )
    add_textbox(
        slide,
        3.25,
        5.65,
        6.9,
        0.35,
        [{"text": f"{applicant.name} | 配置开发 {applicant.target_level} 申报答辩", "size": 11.5, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}],
        margin=0.0,
    )


def build_presentation() -> Path:
    template_path = find_single_file(".pptx", "模版")
    application_path = find_single_file(".xlsx", "申报表")
    standard_path = find_single_file(".xlsx", "资格标准")

    applicant = parse_applicant_data(application_path)
    standards = parse_standard_data(standard_path)

    output_path = ROOT / f"{applicant.name}-{applicant.target_level}专业任职资格答辩PPT.pptx"
    if output_path.exists():
        output_path.unlink()

    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)
    clear_slides(prs)

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    build_cover(cover, applicant)

    contents = prs.slides.add_slide(prs.slide_layouts[1])
    build_contents_slide(contents)

    section_1 = prs.slides.add_slide(prs.slide_layouts[2])
    build_section_slide(section_1, "个人概览与申报结论", "从项目产出、交付质量与标准对标说明本次 DE2 申报基础")

    overview = prs.slides.add_slide(prs.slide_layouts[3])
    build_overview_slide(overview, applicant)

    standard_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_standard_slide(standard_slide, applicant, standards)

    section_2 = prs.slides.add_slide(prs.slide_layouts[5])
    build_section_slide(section_2, "项目产出与交付质量", "覆盖多行业业务场景，形成稳定的独立交付与质量表现")

    output_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_project_output_slide(output_slide, applicant)

    quality_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_project_quality_slide(quality_slide, applicant)

    section_3 = prs.slides.add_slide(prs.slide_layouts[6])
    build_section_slide(section_3, "专业回馈与能力对标", "围绕带教分享、关键能力与流程协同展示我的成长与价值")

    feedback_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_feedback_slide(feedback_slide, applicant)

    capability_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_capability_slide(capability_slide, applicant, standards)

    section_4 = prs.slides.add_slide(prs.slide_layouts[7])
    build_section_slide(section_4, "总结与后续计划", "基于本次申报，持续扩大交付影响力与复用沉淀")

    summary_slide = prs.slides.add_slide(prs.slide_layouts[3])
    build_summary_slide(summary_slide, applicant)

    end_slide = prs.slides.add_slide(prs.slide_layouts[9])
    build_end_slide(end_slide, applicant)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    result = build_presentation()
    print(result)
